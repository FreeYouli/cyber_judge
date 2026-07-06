"""赛博判官答辩PPT生成器 - 在用户笔记本上运行"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

BG = RGBColor(0x0F, 0x1A, 0x2E)
G = RGBColor(0xD4, 0xAF, 0x37)  # gold
W = RGBColor(0xE0, 0xE8, 0xF0)  # white
B = RGBColor(0x7A, 0xB0, 0xD9)  # blue
M = RGBColor(0x5A, 0x7A, 0x9A)  # muted
D = RGBColor(0x1A, 0x2A, 0x4A)  # dark
R = RGBColor(0xFF, 0x6B, 0x6B)  # red

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

def new_slide():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = BG
    return s

def tx(s, l, t, w, h, text, sz=18, clr=W, bd=False, al=PP_ALIGN.LEFT):
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text
    p.font.size = Pt(sz); p.font.color.rgb = clr
    p.font.bold = bd; p.alignment = al

def box(s, l, t, w, h, fc=D):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(l), Inches(t), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = fc
    sh.line.color.rgb = G; sh.line.width = Pt(1.5)

# S1-Cover
s = new_slide()
tx(s, 1, 1.5, 11, 0.8, "赛博判官", 54, G, True, PP_ALIGN.CENTER)
tx(s, 1, 2.7, 11, 0.5, "AI纠纷调解Agent", 22, W, False, PP_ALIGN.CENTER)
sh = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5), Inches(3.5), Inches(3.3), Inches(0.02))
sh.fill.solid(); sh.fill.fore_color.rgb = G; sh.line.fill.background()
tx(s, 1, 4.0, 11, 0.4, "《自然语言处理综合应用实践》专周设计", 14, B, False, PP_ALIGN.CENTER)
tx(s, 1, 4.5, 11, 0.4, "XXX / XXX / XXX", 12, M, False, PP_ALIGN.CENTER)

for n in range(2, 13):
    s = new_slide()
    accent = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(0.8), Inches(0.5), Inches(0.06), Inches(0.5))
    accent.fill.solid(); accent.fill.fore_color.rgb = G; accent.line.fill.background()
    
    if n == 2: tx(s, 1, 0.3, 11, 0.5, "什么是赛博判官", 32, G, True)
    elif n == 3: tx(s, 1, 0.3, 11, 0.5, "解决什么问题", 32, G, True)
    elif n == 4: tx(s, 1, 0.3, 11, 0.5, "系统架构", 32, G, True)
    elif n == 5: tx(s, 1, 0.3, 11, 0.5, "技术栈", 32, G, True)
    elif n == 6: tx(s, 1, 0.3, 11, 0.5, "界面展示", 32, G, True)
    elif n == 7: tx(s, 1, 0.3, 11, 0.5, "案例演示", 32, G, True)
    elif n == 8: tx(s, 1, 0.3, 11, 0.5, "逻辑漏洞审计", 32, G, True)
    elif n == 9: tx(s, 1, 0.3, 11, 0.5, "多维雷达图", 32, G, True)
    elif n == 10: tx(s, 1, 0.3, 11, 0.5, "踩坑与填坑", 32, G, True)
    elif n == 11: tx(s, 1, 0.3, 11, 0.5, "总结与展望", 32, G, True)
    elif n == 12: tx(s, 1, 0.3, 11, 0.5, "谢谢！Q&A", 32, G, True)

prs.save("赛博判官_答辩PPT.pptx")
print("生成完成：赛博判官_答辩PPT.pptx")
